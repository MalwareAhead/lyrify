import tkinter as tk
from tkinter import ttk
import os
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import tkinter.messagebox as messagebox
from tkinter import filedialog
from tkinterdnd2 import DND_FILES, TkinterDnD


class LyricsToPPTXConverter:
    def __init__(self):
        self.window = TkinterDnD.Tk()
        self.window.title("Lyrics to PPTX Converter")
        self.window.geometry("600x400")

        self.output_dir = "presentations"
        if not os.path.exists(self.output_dir):
            os.makedirs(self.output_dir)

        self.setup_ui()

    def setup_ui(self):
        self.drop_frame = ttk.LabelFrame(self.window, text="Drop Files Here")
        self.drop_frame.pack(padx=20, pady=20, fill=tk.BOTH, expand=True)

        self.label = ttk.Label(self.drop_frame,
                               text="Drop .txt files here or click to select files")
        self.label.pack(padx=20, pady=50)

        self.drop_frame.drop_target_register(DND_FILES)
        self.drop_frame.dnd_bind('<<Drop>>', self.handle_drop)
        self.drop_frame.bind("<Button-1>", self.select_files)

    def select_files(self, event=None):
        files = filedialog.askopenfilenames(filetypes=[("Text Files", "*.txt")])
        if files:
            self.process_files(files)

    def handle_drop(self, event):
        files = [f.strip('{}') for f in event.data.split()]
        text_files = [f for f in files if f.lower().endswith('.txt')]
        if text_files:
            self.process_files(text_files)

    def process_files(self, files):
        for file_path in files:
            try:
                self.create_presentation(file_path)
                messagebox.showinfo("Success",
                                    f"Presentation created for {os.path.basename(file_path)}")
            except Exception as e:
                messagebox.showerror("Error",
                                     f"Error processing {os.path.basename(file_path)}: {str(e)}")

    def create_presentation(self, txt_file):
        with open(txt_file, 'r', encoding='utf-8') as file:
            content = file.read()

        slides_content = [slide.strip() for slide in content.split('\n\n')]

        prs = Presentation()
        prs.slide_width = Pt(1920)
        prs.slide_height = Pt(1080)

        slide_height = 1080
        main_text_height = int(slide_height * (2 / 3))

        for i, slide_text in enumerate(slides_content):
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0, 0, 0)

            # Calculate top margin based on number of lines
            line_count = len(slide_text.split('\n'))
            top_margin = 175
            if line_count == 2:
                top_margin = 325
            elif line_count == 3:
                top_margin = 250

            main_text_box = slide.shapes.add_textbox(
                left=Pt(0),
                top=Pt(top_margin),
                width=Pt(1920),
                height=Pt(main_text_height)
            )

            main_text_frame = main_text_box.text_frame
            main_text_frame.text = slide_text
            main_text_frame.vertical_anchor = 1

            for paragraph in main_text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                paragraph.font.name = 'Arial'
                paragraph.font.size = Pt(90)
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
                paragraph.line_spacing = 1.5

            if i < len(slides_content) - 1:
                preview_text_box = slide.shapes.add_textbox(
                    left=Pt(0),
                    top=Pt(900),
                    width=Pt(1920),
                    height=Pt(100)
                )

                preview_text_frame = preview_text_box.text_frame
                next_slide_first_line = slides_content[i + 1].split('\n')[0]
                preview_text_frame.text = next_slide_first_line

                preview_paragraph = preview_text_frame.paragraphs[0]
                preview_paragraph.alignment = PP_ALIGN.CENTER
                preview_paragraph.font.name = 'Arial'
                preview_paragraph.font.size = Pt(45)
                preview_paragraph.font.color.rgb = RGBColor(128, 128, 128)

        # Add final empty slide
        last_slide = prs.slides.add_slide(prs.slide_layouts[6])
        last_slide.background.fill.solid()
        last_slide.background.fill.fore_color.rgb = RGBColor(0, 0, 0)

        output_filename = os.path.splitext(os.path.basename(txt_file))[0] + '.pptx'
        output_path = os.path.join(self.output_dir, output_filename)
        prs.save(output_path)

    def run(self):
        self.window.mainloop()


if __name__ == "__main__":
    app = LyricsToPPTXConverter()
    app.run()